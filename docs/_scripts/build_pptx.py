# -*- coding: utf-8 -*-
"""Genera la presentacion de sustentacion ReciclaPe (.pptx) 16:9."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
IMG = os.path.join(BASE, 'Anexo-A-Diagramas', 'img')

GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK = RGBColor(0x1B, 0x5E, 0x20)
GREY = RGBColor(0x42, 0x42, 0x42)
LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h, text, size=18, color=GREY, bold=False,
            align=PP_ALIGN.LEFT, italic=False, font='Arial'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=GREY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2); r.font.color.rgb = color
        r.font.name = 'Arial'
        p.space_after = Pt(6)
    return tb


def header(slide, title):
    rect(slide, 0, 0, SW, Inches(1.1), GREEN)
    textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8),
            title, size=26, color=WHITE, bold=True)
    # franja inferior
    rect(slide, 0, SH - Inches(0.25), SW, Inches(0.25), DARK)


def content_image(slide, path, top=Inches(1.3), maxw=Inches(8.0), maxh=Inches(5.6), left=None):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(maxw / iw, maxh / ih)
    w = Emu(int(iw * ratio)); h = Emu(int(ih * ratio))
    if left is None:
        left = Emu(int((SW - w) / 2))
    slide.shapes.add_picture(path, left, top, width=w, height=h)


# ---------------------------------------------------------------- 1 Portada
s = add_slide()
rect(s, 0, 0, SW, SH, GREEN)
rect(s, 0, Inches(2.6), SW, Inches(2.3), WHITE)
textbox(s, Inches(1), Inches(2.85), Inches(11.3), Inches(1.1),
        "ReciclaPe", size=54, color=DARK, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.95), Inches(11.3), Inches(0.7),
        "Plataforma de gestión de reciclaje vecinal basada en microservicios",
        size=20, color=GREY, align=PP_ALIGN.CENTER, italic=True)
textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
        "Desarrollo de Aplicaciones Web II (4697)  |  CIBERTEC  |  Ciclo Sexto  |  2026",
        size=16, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
textbox(s, Inches(1), Inches(5.9), Inches(11.3), Inches(1.0),
        "Equipo: [Coordinador y 4 integrantes]   ·   Profesor: [Nombre del profesor]",
        size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 2 Agenda
s = add_slide(); header(s, "Agenda")
bullets(s, Inches(1.2), Inches(1.5), Inches(11), Inches(5.5), [
    ("El problema: residuos sólidos y reciclaje en Lima", 0),
    ("La solución ReciclaPe y sus objetivos", 0),
    ("Arquitectura de microservicios", 0),
    ("Modelo de datos y API REST", 0),
    ("Seguridad (JWT + BCrypt) y flujo de recojo", 0),
    ("Frontend Angular por roles", 0),
    ("Despliegue (Docker) y observabilidad", 0),
    ("Cumplimiento de rúbrica y sílabo", 0),
    ("Demo y conclusiones", 0),
], size=22)

# ---------------------------------------------------------------- 3 Problema
s = add_slide(); header(s, "El problema (Análisis SEPTE)")
bullets(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), [
    ("Lima genera 9 343 toneladas/día de residuos sólidos municipales.", 0),
    ("Solo el 3.9% de los reciclables se recupera formalmente.", 0),
    ("~108 000 recicladores en el Perú; solo ~8% formalizados.", 0),
    ("Menos del 15% de los hogares limeños segrega en origen.", 0),
    ("Marco legal exige a las municipalidades programas de segregación (Ley 27314, 29419).", 0),
    ("Oportunidad: reducir la fricción entre vecino y reciclador y medir el impacto.", 0),
], size=20)

# ---------------------------------------------------------------- 4 Solucion
s = add_slide(); header(s, "La solución: ReciclaPe")
bullets(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), [
    ("Plataforma web que conecta vecinos con recicladores formales.", 0),
    ("Programa y da seguimiento a recojos de residuos segregados.", 0),
    ("Calcula el impacto ambiental: kg recuperados y CO₂ evitado.", 0),
    ("Tres roles: Vecino, Reciclador y Administrador municipal.", 0),
    ("Backend de microservicios Spring Boot + frontend Angular.", 0),
], size=22)

# ---------------------------------------------------------------- 5 Objetivos
s = add_slide(); header(s, "Objetivos (SMART)")
bullets(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.7), [
    ("OBJ 1 — Backend de ≥3 microservicios (auth, recojos, reportes) con JWT+BCrypt, "
     "comunicación Feign y RabbitMQ, ≥20 endpoints REST tras un API Gateway.", 0),
    ("OBJ 2 — Frontend Angular por rol que consume el 100% de los endpoints y muestra "
     "un dashboard de impacto con ≥3 indicadores.", 0),
    ("OBJ 3 — Contenerizar todo con Docker Compose y publicar un dashboard de "
     "observabilidad (Actuator + Prometheus + Grafana).", 0),
], size=20)

# ---------------------------------------------------------------- 6 Arquitectura
s = add_slide(); header(s, "Arquitectura de microservicios")
content_image(s, os.path.join(IMG, '03-arquitectura-general.png'))

# ---------------------------------------------------------------- 7 Modelo datos
s = add_slide(); header(s, "Modelo de datos")
content_image(s, os.path.join(IMG, '02-modelo-entidad-relacion.png'), maxh=Inches(5.7))

# ---------------------------------------------------------------- 8 Seguridad
s = add_slide(); header(s, "Seguridad: login con JWT y BCrypt")
content_image(s, os.path.join(IMG, '05-secuencia-login.png'), maxw=Inches(11.5))

# ---------------------------------------------------------------- 9 Flujo recojo
s = add_slide(); header(s, "Flujo completo de recojo")
content_image(s, os.path.join(IMG, '06-secuencia-recojo.png'), maxw=Inches(11.5))

# ---------------------------------------------------------------- 10 Endpoints
s = add_slide(); header(s, "API REST (GET / POST / PUT / DELETE)")
content_image(s, os.path.join(IMG, '09-mapa-endpoints.png'), maxw=Inches(12.2))

# ---------------------------------------------------------------- 11 Frontend
s = add_slide(); header(s, "Frontend Angular por roles")
content_image(s, os.path.join(IMG, '04-componentes-angular.png'), maxw=Inches(11.5))

# ---------------------------------------------------------------- 12 Despliegue
s = add_slide(); header(s, "Despliegue y observabilidad")
content_image(s, os.path.join(IMG, '08-despliegue.png'), maxw=Inches(8.2), left=Inches(0.4))
bullets(s, Inches(8.8), Inches(1.5), Inches(4.2), Inches(5.3), [
    ("docker compose up levanta todo el stack.", 0),
    ("Actuator expone /prometheus.", 0),
    ("Prometheus recolecta métricas.", 0),
    ("Grafana muestra latencia y errores.", 0),
    ("RabbitMQ para eventos asíncronos.", 0),
], size=16)

# ---------------------------------------------------------------- 13 Cumplimiento
s = add_slide(); header(s, "Cumplimiento: rúbrica + sílabo")
bullets(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.7), [
    ("Rúbrica: login REST + BCrypt, CRUD GET/POST/PUT/DELETE, Angular consume todo, "
     "pruebas de capa de datos (JUnit @DataJpaTest).", 0),
    ("Sílabo U1: Feign, RabbitMQ, Spring Cloud (Gateway, Eureka).", 0),
    ("Sílabo U2: JWT, Spring Security, Resilience4J, Docker.", 0),
    ("Sílabo U3: Actuator, Prometheus, Grafana (observabilidad).", 0),
    ("Parte formal: SEPTE con fuentes citadas, objetivos SMART, justificación, conclusiones.", 0),
], size=19)

# ---------------------------------------------------------------- 14 Cronograma
s = add_slide(); header(s, "Hitos del proyecto")
bullets(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.5), [
    ("Semana 5 — T1 (laboratorio individual).", 0),
    ("Semana 9 — T2 (laboratorio individual).", 0),
    ("Semana 11 — AP1: auth + recojo + Angular login + 1 CRUD + tests + compose.", 0),
    ("Semana 13-14 — EF/SP1: todos los microservicios + Gateway + Feign + RabbitMQ + "
     "Actuator/Prometheus + Angular completo + video.", 0),
], size=20)

# ---------------------------------------------------------------- 15 Conclusiones
s = add_slide(); header(s, "Conclusiones")
bullets(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), [
    ("La arquitectura de microservicios alinea el proyecto con el logro del curso y separa "
     "responsabilidades con claridad.", 0),
    ("El registro digital de recojos y el cálculo de impacto dan trazabilidad útil para las "
     "municipalidades (Ley 27314, Programa de Incentivos del MEF).", 0),
    ("Integrar seguridad, resiliencia y observabilidad desde temprano es factible a escala "
     "académica y transferible a lo profesional.", 0),
], size=20)

# ---------------------------------------------------------------- 16 Gracias
s = add_slide()
rect(s, 0, 0, SW, SH, GREEN)
textbox(s, Inches(1), Inches(2.9), Inches(11.3), Inches(1.2),
        "¡Gracias!", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
        "ReciclaPe — Desarrollo de Aplicaciones Web II (4697)",
        size=20, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)

out = os.path.join(BASE, 'Presentacion-Sustentacion.pptx')
prs.save(out)
print("Guardado:", out, "-", len(prs.slides._sldIdLst), "slides")
